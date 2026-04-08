from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()

    # Slide 1: About the project
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "1. About the Project"
    
    tf = content.text_frame
    tf.text = "The 'OutPass Monitor' is a web-based application designed to manage student and parent movements in and out of the college."
    p = tf.add_paragraph()
    p.text = "It replaces manual registers with an automated online system."
    p = tf.add_paragraph()
    p.text = "The project streamlines the approval process involving Coordinators, HODs, and Gate security."
    p = tf.add_paragraph()
    p.text = "It ensures secure, tracked, and efficient management of out-passes."

    # Slide 2: Hardware & Software Specifications
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "2. Hardware & Software Specifications"
    
    tf = content.text_frame
    tf.text = "Hardware Specifications:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "RAM: 2GB"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "ROM: 256GB"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Processor: Windows 7 32-bit compatible processor"
    
    p = tf.add_paragraph()
    p.text = "Software Specifications:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "OS: Windows 7"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "PHP Version: 7.3.2"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "MySQL Version: 5.10"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "XAMPP Version: 3.2.2"

    # Slide 3: Problem Definition
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "3. Problem Definition"
    
    tf = content.text_frame
    tf.text = "The existing manual system has several drawbacks:"
    p = tf.add_paragraph()
    p.text = "It is slow and time-consuming for students to get physical signatures."
    p = tf.add_paragraph()
    p.text = "Paper records are hard to maintain and can be easily lost or damaged."
    p = tf.add_paragraph()
    p.text = "There is no real-time tracking of who is inside or outside the campus."
    p = tf.add_paragraph()
    p.text = "Communication between departments and the gate is difficult."

    # Slide 4: System Study
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "4. System Study"
    
    tf = content.text_frame
    tf.text = "Existing System:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Manual paper-based applications."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Physical movement required for approvals."
    
    p = tf.add_paragraph()
    p.text = "Disadvantages:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Time wastage and lack of data security."
    
    p = tf.add_paragraph()
    p.text = "Need for New System:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "To automate the workflow and ensure data integrity."

    # Slide 5: Proposed System
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "5. Proposed System"
    
    tf = content.text_frame
    tf.text = "The proposed 'OutPass Monitor' is a web application."
    p = tf.add_paragraph()
    p.text = "Users (Students/Parents) can log in and submit requests online."
    p = tf.add_paragraph()
    p.text = "Coordinators and HODs can view and approve requests from their dashboard."
    p = tf.add_paragraph()
    p.text = "The status is updated in real-time for the student."
    p = tf.add_paragraph()
    p.text = "The Gatekeeper can verify the approved pass digitally."

    # Slide 6: Dataflow Diagram
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "6. Dataflow Diagram"
    
    tf = content.text_frame
    tf.text = "Flow of Data in the System:"
    p = tf.add_paragraph()
    p.text = "1. Student registers and logs in."
    p = tf.add_paragraph()
    p.text = "2. Student submits an Out-Pass Request."
    p = tf.add_paragraph()
    p.text = "3. Request data is stored in the MySQL Database."
    p = tf.add_paragraph()
    p.text = "4. Coordinator views and approves the request."
    p = tf.add_paragraph()
    p.text = "5. HOD views and gives final approval."
    p = tf.add_paragraph()
    p.text = "6. Gate verifies the status and allows exit."

    prs.save('GPM_Project_Presentation.pptx')
    print("Presentation created successfully: GPM_Project_Presentation.pptx")

if __name__ == "__main__":
    create_ppt()
