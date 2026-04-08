from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


def wipe_slide(slide):
    # Remove all shapes to start with a clean slate if needed, 
    # but since we are creating new slides, we just control placement.
    pass

def apply_title_style(shape, text):
    shape.text = text
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.font.name = 'Arial Black'
    paragraph.font.size = Pt(32)
    paragraph.font.color.rgb = RGBColor(44, 62, 80) # Dark Blue-Grey
    paragraph.alignment = PP_ALIGN.LEFT

def apply_text_style(paragraph, font_size=20, bold=False, color=None):
    paragraph.font.name = 'Segoe UI'
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    if color:
        paragraph.font.color.rgb = color
    else:
        paragraph.font.color.rgb = RGBColor(60, 60, 60) # Dark Gray
    paragraph.space_after = Pt(12)

def add_header_footer(slide):
    # Header Bar
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(0.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(52, 152, 219) # Bright Blue
    shape.line.fill.background()
    
    # Footer Bar
    top = Inches(7.1)
    height = Inches(0.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(44, 62, 80) # Dark Blue-Grey
    shape.line.fill.background()

def create_ppt():
    prs = Presentation()

    # --- Slide 1: About the Project ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    add_header_footer(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    apply_title_style(title_box, "1. About the Project")
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "The 'OutPass Monitor' is a comprehensive web-based application designed to modernize the management of student and parent movements in and out of the college campus."
    apply_text_style(p, 22)
    
    bullets = [
        "Replaces traditional manual registers with a digital system.",
        "Streamlines approval workflows between Coordinators, HODs, and Gate Security.",
        "Ensures secure, real-time tracking of student exits and entries.",
        "Enhances administrative efficiency and data accuracy."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        apply_text_style(p, 20)

    # --- Slide 2: Hardware & Software Specifications ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    apply_title_style(title_box, "2. Hardware & Software Specifications")
    
    # Left Column: Hardware
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.25), Inches(4.5))
    tf = left_box.text_frame
    
    p = tf.add_paragraph()
    p.text = "Hardware Specs"
    apply_text_style(p, 22, bold=True, color=RGBColor(41, 128, 185))
    
    hw_specs = [
        ("RAM", "2GB"),
        ("ROM", "256GB"),
        ("Processor", "Windows 7 32-bit compatible")
    ]
    for label, val in hw_specs:
        p = tf.add_paragraph()
        p.text = f"{label}: {val}"
        apply_text_style(p, 18)

    # Right Column: Software
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(2), Inches(4.25), Inches(4.5))
    tf = right_box.text_frame
    
    p = tf.add_paragraph()
    p.text = "Software Specs"
    apply_text_style(p, 22, bold=True, color=RGBColor(41, 128, 185))
    
    sw_specs = [
        ("OS", "Windows 7"),
        ("PHP Version", "7.3.2"),
        ("MySQL Version", "5.10"),
        ("XAMPP Version", "3.2.2")
    ]
    for label, val in sw_specs:
        p = tf.add_paragraph()
        p.text = f"{label}: {val}"
        apply_text_style(p, 18)

    # --- Slide 3: Problem Definition ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    apply_title_style(title_box, "3. Problem Definition")
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))
    tf = content_box.text_frame
    
    p = tf.add_paragraph()
    p.text = "Limitations of the Manual System:"
    apply_text_style(p, 22, bold=True, color=RGBColor(192, 57, 43))
    
    problems = [
        "Slow Process: Obtaining physical signatures from multiple authorities is time-consuming.",
        "Data Loss Risk: Paper records are prone to damage, misplacement, and wear over time.",
        "Lack of Visibility: No real-time tracking of student location (inside/outside).",
        "Communication Gaps: Poor coordination between departments and the main gate."
    ]
    for item in problems:
        p = tf.add_paragraph()
        p.text = "• " + item
        apply_text_style(p, 20)

    # --- Slide 4: System Study ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    apply_title_style(title_box, "4. System Study")
    
    # Comparison layout
    # Existing
    box_1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.25), Inches(4))
    tf = box_1.text_frame
    p = tf.add_paragraph()
    p.text = "Existing System"
    apply_text_style(p, 22, bold=True, color=RGBColor(127, 140, 141)) # Grey
    
    existing_pts = [
        "Paper-based applications",
        "Physical movement required",
        "High manual effort",
        "Prone to errors"
    ]
    for pt in existing_pts:
        p = tf.add_paragraph()
        p.text = "- " + pt
        apply_text_style(p, 18)

    # Proposed
    box_2 = slide.shapes.add_textbox(Inches(5.0), Inches(2), Inches(4.5), Inches(4))
    tf = box_2.text_frame
    p = tf.add_paragraph()
    p.text = "Proposed System"
    apply_text_style(p, 22, bold=True, color=RGBColor(39, 174, 96)) # Green
    
    proposed_pts = [
        "Web-based application dashboard",
        "Remote/Online approval",
        "Automated workflows",
        "Secure digital records"
    ]
    for pt in proposed_pts:
        p = tf.add_paragraph()
        p.text = "✓ " + pt
        apply_text_style(p, 18)

    # --- Slide 5: Proposed System ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    apply_title_style(title_box, "5. Proposed System")
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))
    tf = content_box.text_frame
    
    p = tf.add_paragraph()
    p.text = "Key Features & Workflow:"
    apply_text_style(p, 22, bold=True, color=RGBColor(41, 128, 185))
    
    features = [
        "Role-Based Access: Separate logins for Students, Parents, Coordinators, HODs, and Gate.",
        "Online Requests: Students/Parents submit leave or out-pass requests digitally.",
        "Multi-Level Approval: Requests are automatically routed to the correct authority.",
        "Real-Time Status: Students can check approval status instantly.",
        "Gate Verification: Security guards verify approved passes via the system before allowing exit."
    ]
    for f in features:
        p = tf.add_paragraph()
        p.text = "• " + f
        apply_text_style(p, 20)

    # --- Slide 6: Dataflow Diagram ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    apply_title_style(title_box, "6. Dataflow Diagram")
    
    # --- Slide 6: Dataflow Diagram ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    apply_title_style(title_box, "6. Dataflow Diagram")
    
    # DFD Drawing Functions
    def add_entity(text, left, top, color=RGBColor(52, 152, 219)):
        width = Inches(1.5)
        height = Inches(0.8)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(44, 62, 80)
        shape.text = text
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.size = Pt(14)
        return shape

    def add_process(text, left, top, color=RGBColor(46, 204, 113)):
        width = Inches(1.5)
        height = Inches(0.8)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(44, 62, 80)
        shape.text = text
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.size = Pt(14)
        return shape

    def add_datastore(text, left, top, color=RGBColor(241, 196, 15)):
        width = Inches(1.5)
        height = Inches(0.8)
        # Flowchart Manual Operation or similar looks like a datastore open ended
        # Or just use Can
        shape = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_MAGNETIC_DISK, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(44, 62, 80)
        shape.text = text
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.size = Pt(12)
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
        return shape

    def connect(shape1, shape2, text=""):
        # Simple midpoint connection
        x1 = shape1.left + shape1.width
        y1 = shape1.top + shape1.height/2
        x2 = shape2.left
        y2 = shape2.top + shape2.height/2
        
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.width = Pt(2)
        line.line.fill.solid()
        line.line.fill.fore_color.rgb = RGBColor(127, 140, 141)

        
        if text:
            lbl = slide.shapes.add_textbox(x1 + (x2-x1)/2 - Inches(0.5), y1 - Inches(0.3), Inches(1), Inches(0.3))
            lbl.text_frame.text = text
            lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            lbl.text_frame.paragraphs[0].font.size = Pt(10)

    # Layout:
    # [Student] --> (Login) --> [DB]
    #                             |
    # [Parent]  --> (Login) ------+
    #                             |
    #                           (Check)
    #                             |
    #                        [Coordinator] --> [HOD] --> [Gate]
    
    # Row 1: Users
    y_users = Inches(2.0)
    student = add_entity("Student", Inches(0.5), y_users)
    parent = add_entity("Parent", Inches(0.5), y_users + Inches(1.2))
    
    # Process Column
    y_proc = Inches(2.5)
    login_proc = add_process("Login &\nRequest", Inches(2.5), y_proc)
    
    # DB Column
    db = add_datastore("MySQL\nDatabase", Inches(4.5), y_proc)
    
    # Approval Flow
    y_approve = Inches(4.5)
    coord = add_entity("Coordinator", Inches(2.5), y_approve)
    hod = add_entity("HOD", Inches(4.5), y_approve)
    gate = add_entity("Gate", Inches(6.5), y_approve)
    
    # Connections
    connect(student, login_proc)
    connect(parent, login_proc)
    connect(login_proc, db, "Save")
    
    # Connect DB to Coord (Logic: Coord checks DB)
    # Drawing vertical arrow manually for DB -> Coord
    v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, db.left + db.width/2, db.top + db.height, coord.left + coord.width/2, coord.top)
    v_line.line.width = Pt(2)
    v_line.line.fill.solid()
    v_line.line.fill.fore_color.rgb = RGBColor(127, 140, 141)

    
    # Approval Chain
    connect(coord, hod, "Approve")
    connect(hod, gate, "Finalize")
    
    # Final Outcome
    exit_node = add_process("Exit\nAllowed", Inches(8.5), y_approve, RGBColor(231, 76, 60))
    connect(gate, exit_node, "Verify")

    prs.save('GPM_Project_Presentation_Pro_v5.pptx')
    print("Presentation updated successfully: GPM_Project_Presentation_Pro_v5.pptx")

if __name__ == "__main__":
    create_ppt()
